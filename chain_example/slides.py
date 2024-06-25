from pptx import Presentation
from pptx.util import Inches, Pt

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Using the bullet layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content[0]
    
    for line in content[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 1
    
    return slide

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Unveiling the Mischievous World of Duendes"
subtitle.text = "Folklore's Tiny Tricksters"

# Slide 2: Introduction
create_slide(prs, "Introduction", [
    "Step into a realm where pint-sized pranksters rule and mischief reigns supreme!",
    "Duendes are mythical creatures in Latin American and Filipino folklore",
    "Known for their small size and mischievous nature",
    "Often associated with households and natural environments"
])

# Slide 3: Characteristics of Duendes
create_slide(prs, "Characteristics of Duendes", [
    "Physical appearance and behavior:",
    "Usually depicted as small, humanoid creatures",
    "Often wear pointed hats or unusual clothing",
    "Known for their playful and sometimes troublesome antics",
    "Said to possess magical or supernatural abilities"
])

# Slide 4: Cultural Significance
create_slide(prs, "Cultural Significance", [
    "Importance in folklore and traditions:",
    "Feature prominently in stories and legends across various cultures",
    "Often used to explain mysterious occurrences or missing objects",
    "Some believe they can bring good or bad luck to households",
    "Influence on art, literature, and popular culture"
])

# Slide 5: Modern Interpretations
create_slide(prs, "Modern Interpretations", [
    "Contemporary views on Duendes:",
    "Continued presence in modern storytelling and media",
    "Subject of paranormal investigations and folklore studies",
    "Symbolic representation of the mysterious and unexplained",
    "Enduring fascination with these mythical creatures"
])

# Save the presentation
prs.save('Duendes_Presentation.pptx')
print("Presentation created successfully!")